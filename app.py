import os
import time
import re
import pathlib

import requests
import openai
from embedchain import App
from serpapi import GoogleSearch
from pptx import Presentation
from pptx.util import Inches

from pptx import Presentation
from pptx.util import Inches, Pt
import gradio as gr

import torch

from PIL import Image
import qrcode
from pathlib import Path
from multiprocessing import cpu_count
import requests
import io
import os
from PIL import Image


from diffusers import (
    StableDiffusionControlNetPipeline,
    ControlNetModel,
    DDIMScheduler,
    DPMSolverMultistepScheduler,
    DEISMultistepScheduler,
    HeunDiscreteScheduler,
    EulerDiscreteScheduler,
    EulerAncestralDiscreteScheduler,
)


openai.api_key = os.environ['OPENAI_API_KEY']
def gpt(user_prompt: str) -> str:
    response = openai.Completion.create(
      model="text-davinci-003",
      prompt=user_prompt,
      temperature=0,
      max_tokens=200,
      top_p=1,
      frequency_penalty=0,
      presence_penalty=0)
    return response["choices"][0]["text"]

def get_results(query:str, topic:str,index=0)->list[str]:
  combined_q = gpt(f'combine these "{query}" + "{topic}" words  and generate one heading')
  print(f'{query = }, {topic = }, {combined_q = }')
  
  try:
    params = {
    "engine": "google",
    "q": combined_q,
    "api_key": os.environ[f'SERPAPI_API_KEY{index}']
    }
    search = GoogleSearch(params)
    results = search.get_dict()
  except Exception as e:
    print(e)
    get_results(query, topic,index=index+1)



  organic_results = results["organic_results"]
  return organic_results

def extract_points(query:str, topic:str)->list[str]:
  # print('--Sleep--')
  time.sleep(60)
  organic_results = get_results(query, topic)
  embd_chain = App()
  for index, dct in enumerate(organic_results):
    try:
      embd_chain.add('web_page',dct['link'])
    except requests.exceptions.SSLError:
      continue
    except openai.error.RateLimitError:
      break
  print('--sleep--')
  time.sleep(60)
  embd_chain_q = embd_chain.query(f'highlight 7 important points')

  return 
# Add the title slide

def add_slide(prs, title, content, title_font_size=Pt(36), content_font_size=Pt(18)):
    slide_layout = prs.slide_layouts[1]  # Use the layout for "Title and Content"
    slide = prs.slides.add_slide(slide_layout)

    # Set the title and content text
    slide.shapes.title.text = title
    text_box = slide.placeholders[1]
    text_box.text = content

    # Change the font size for title and content text
    title_text_frame = slide.shapes.title.text_frame
    content_text_frame = text_box.text_frame
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = title_font_size

    for paragraph in content_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = content_font_size


def add_title_slide(prs, title, title_font_size=Pt(44)):
    slide_layout = prs.slide_layouts[0]  # Use the layout for "Title Slide"
    slide = prs.slides.add_slide(slide_layout)

    # Set the title and subtitle text
    slide.shapes.title.text = title


    # Change the font size for title and subtitle text
    title_text_frame = slide.shapes.title.text_frame
    
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = title_font_size


def main(user_query:str)->dict[str, str]:
  res = gpt(f'You are assisting me in creating a presentation on "{user_query}" Please generate 5 informative side headings for the slides. Each heading should be concise and reflect a key aspect of the topic.')
  topics = re.sub(r'[\d.]','',res.strip()).split('\n')
  print(f'{topics = }')
  ppt_points = { topic: extract_points(topic, user_query)
            for topic in topics}
  prs = Presentation()
  add_title_slide(prs,user_query, title_font_size=Pt(44))

  # Data for content slides
  
  # Adding each key-value pair as a slide in the presentation with custom font sizes
  for key, value in ppt_points.items():
      add_slide(prs, key, value, title_font_size=Pt(36), content_font_size=Pt(18))

  # Save the presentation
  prs.save(f'{user_query}.pptx')

  return f'{user_query}.pptx'

controlnet = ControlNetModel.from_pretrained(
    "monster-labs/control_v1p_sd15_qrcode_monster", 
    torch_dtype=torch.float16

).to('cpu')

pipe = StableDiffusionControlNetPipeline.from_pretrained(
    "runwayml/stable-diffusion-v1-5",
    controlnet=controlnet,
    safety_checker=None,
    torch_dtype=torch.float16
    
    
).to('cuda')
pipe.enable_xformers_memory_efficient_attention()


SAMPLER_MAP = {
    "DPM++ Karras SDE": lambda config: DPMSolverMultistepScheduler.from_config(config, use_karras=True, algorithm_type="sde-dpmsolver++"),
    "DPM++ Karras": lambda config: DPMSolverMultistepScheduler.from_config(config, use_karras=True),
    "Heun": lambda config: HeunDiscreteScheduler.from_config(config),
    "Euler a": lambda config: EulerAncestralDiscreteScheduler.from_config(config),
    "Euler": lambda config: EulerDiscreteScheduler.from_config(config),
    "DDIM": lambda config: DDIMScheduler.from_config(config),
    "DEIS": lambda config: DEISMultistepScheduler.from_config(config),
}


def create_code(content: str):
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=16,
        border=0,
    )
    qr.add_data(content)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")

    # find smallest image size multiple of 256 that can fit qr
    offset_min = 8 * 16
    w, h = img.size
    w = (w + 255 + offset_min) // 256 * 256
    h = (h + 255 + offset_min) // 256 * 256
    if w > 1024:
        raise gr.Error("QR code is too large, please use a shorter content")
    bg = Image.new('L', (w, h), 128)

    # align on 16px grid
    coords = ((w - img.size[0]) // 2 // 16 * 16,
              (h - img.size[1]) // 2 // 16 * 16)
    bg.paste(img, coords)
    return bg


def inference(
    qr_code_content: str,
    prompt: str,
    negative_prompt: str,
    guidance_scale: float = 10.0,
    controlnet_conditioning_scale: float = 2.0,
    seed: int = -1,
    sampler="Euler a",
):


    pipe.scheduler = SAMPLER_MAP[sampler](pipe.scheduler.config)

    generator = torch.manual_seed(seed) if seed != -1 else torch.Generator()

    print("Generating QR Code from content")
    qrcode_image = create_code(qr_code_content)

    # hack due to gradio examples
    init_image = qrcode_image

    out = pipe(
        prompt=prompt,
        negative_prompt=negative_prompt,
        image=qrcode_image,
        width=qrcode_image.width,
        height=qrcode_image.height,
        guidance_scale=float(guidance_scale),
        controlnet_conditioning_scale=float(controlnet_conditioning_scale),
        
        num_inference_steps=40,
    )
    return out.images[0]

import gradio as gr


with gr.Blocks() as demo:
  with gr.Tab('Presentation'):
    with gr.Row():
      with gr.Column():
        txt = gr.Textbox(label="Your Query")
      with gr.Column():
        file = gr.File()

    btn = gr.Button('Create Presentation')

    
    btn.click(main, txt, file)
  with gr.Tab('Share'):
    gr.Markdown('This feature needs GPU to run')
    with gr.Row():
        with gr.Column():
            qr_code_content = gr.Textbox(
                label="QR Code Content or URL",
                info="The text you want to encode into the QR code",
                value="",
            )

            prompt = gr.Textbox(
                label="Prompt",
                info="Prompt that guides the generation towards",
            )
            negative_prompt = gr.Textbox(
                label="Negative Prompt",
                value="ugly, disfigured, low quality, blurry, nsfw",
                info="Prompt that guides the generation away from",
            )

            with gr.Accordion(
                label="Params: The generated QR Code functionality is largely influenced by the parameters detailed below",
                open=True,
            ):
                controlnet_conditioning_scale = gr.Slider(
                    minimum=0.5,
                    maximum=2.5,
                    step=0.01,
                    value=1.5,
                    label="Controlnet Conditioning Scale",
                    info="""Controls the readability/creativity of the QR code.
                    High values: The generated QR code will be more readable.
                    Low values: The generated QR code will be more creative.
                    """
                )
                guidance_scale = gr.Slider(
                    minimum=0.0,
                    maximum=25.0,
                    step=0.25,
                    value=7,
                    label="Guidance Scale",
                    info="Controls the amount of guidance the text prompt guides the image generation"
                )
                sampler = gr.Dropdown(choices=list(
                    SAMPLER_MAP.keys()), value="Euler a", label="Sampler")
                seed = gr.Number(
                    minimum=-1,
                    maximum=9999999999,
                    step=1,
                    value=2313123,
                    label="Seed",
                    randomize=True,
                    info="Seed for the random number generator. Set to -1 for a random seed"
                )
            with gr.Row():
                run_btn = gr.Button("Run")
        with gr.Column():
            result_image = gr.Image(label="Result Image", elem_id="result_image")
    run_btn.click(
        inference,
        inputs=[
            qr_code_content,
            prompt,
            negative_prompt,
            guidance_scale,
            controlnet_conditioning_scale,
            seed,
            sampler,
        ],
        outputs=[result_image],
    )

    gr.Examples(
        examples=[
            [
                "test",
                "Baroque rococo architecture, architectural photography, post apocalyptic New York, hyperrealism, [roots], hyperrealistic, octane render, cinematic, hyper detailed, 8K",
                "",
                7,
                1.6,
                2592353769,
                "Euler a",
            ],
            [
                "https://qrcodemonster.art",
                "a centered render of an ancient tree covered in bio - organic micro organisms growing in a mystical setting, cinematic, beautifully lit, by tomasz alen kopera and peter mohrbacher and craig mullins, 3d, trending on artstation, octane render, 8k",
                "",
                7,
                1.57,
                259235398,
                "Euler a",
            ],
            [
                "test",
                "3 cups of coffee with coffee beans around",
                "",
                7,
                1.95,
                1889601353,
                "Euler a",
            ],
            [
                "https://huggingface.co",
                "A top view picture of a sandy beach with a sand castle, beautiful lighting, 8k, highly detailed",
                "sky",
                7,
                1.15,
                46200,
                "Euler a",
            ],
            [
                "test",
                "A top view picture of a sandy beach, organic shapes, beautiful lighting, bumps and shadows, 8k, highly detailed",
                "sky, water, squares",
                7,
                1.25,
                46220,
                "Euler a",
            ],
        ],
        fn=inference,
        inputs=[
            qr_code_content,
            prompt,
            negative_prompt,
            guidance_scale,
            controlnet_conditioning_scale,
            seed,
            sampler,
        ],
        outputs=[result_image],
        
    )



demo.launch(debug=True)
